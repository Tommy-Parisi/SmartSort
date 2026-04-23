from pathlib import Path
from ...core.utils import token_cap
from ..identity_utils import clean_filename_stem


def _slide_title(slide) -> str:
    """Return the title placeholder text, or first non-empty text on the slide."""
    try:
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 0:  # 0 = title
                t = shape.text.strip()
                if t:
                    return t
    except Exception:
        pass
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            return shape.text.strip().split("\n")[0]
    return ""


class PptxExtractorAgent:
    def extract(self, file_path: str) -> str:
        try:
            from pptx import Presentation
        except ImportError:
            return clean_filename_stem(file_path)

        try:
            prs = Presentation(file_path)
        except Exception:
            return clean_filename_stem(file_path)

        titles = []
        for slide in prs.slides:
            t = _slide_title(slide)
            if t:
                titles.append(t)

        return token_cap(" ".join(titles)) if titles else clean_filename_stem(file_path)
